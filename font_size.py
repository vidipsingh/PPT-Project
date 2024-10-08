import tkinter as tk
from tkinter import filedialog
from pptx import Presentation

def check_font_sizes_pptx(file_path):
    prs = Presentation(file_path)
    slides_font_sizes = {}

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_font_sizes = {}

        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame is not None:
                for paragraph in shape.text_frame.paragraphs:
                    paragraph_text = "".join([run.text for run in paragraph.runs]).strip()
                    
                    if paragraph_text:  # Only proceed if there's actual text
                        # Assuming the font size of the first run applies to the paragraph
                        font_size = paragraph.runs[0].font.size.pt if paragraph.runs and paragraph.runs[0].font.size else None
                        if font_size:
                            slide_font_sizes[paragraph_text] = font_size

        slides_font_sizes[slide_num] = slide_font_sizes

    return slides_font_sizes

def select_file():
    # Open file dialog and allow the user to select a .pptx file
    root = tk.Tk()
    root.withdraw()  # Hide the main tkinter window
    file_path = filedialog.askopenfilename(
        filetypes=[("PowerPoint files", "*.pptx")],
        title="Select a PowerPoint File"
    )
    return file_path

if __name__ == "__main__":
    file_path = select_file()  # Select the file using a dialog
    if file_path:  # If a file was selected
        print(f"Analyzing file: {file_path.split('/')[-1]}")  # Display the file name
        slides_font_sizes = check_font_sizes_pptx(file_path)
        
        for slide_num, font_sizes in slides_font_sizes.items():
            print(f"\nSlide {slide_num}:")
            for text, size in font_sizes.items():
                print(f"  Text: '{text}', Font Size: {size} pt")
    else:
        print("No file selected.")
