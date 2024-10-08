import tkinter as tk
from tkinter import filedialog
from pptx import Presentation

def get_font_from_theme(prs, font_name):
    """
    Map the theme-based font names to actual font family if applicable.
    """
    # Extract theme fonts (e.g., Body and Heading fonts)
    if prs.core_properties:
        theme_fonts = prs.slide_master.text_styles
        if font_name == "Body":
            return theme_fonts.body_font.font.name
        elif font_name == "Heading":
            return theme_fonts.heading_font.font.name

    return font_name

def check_font_theme_pptx(file_path):
    prs = Presentation(file_path)
    slides_font_themes = {}

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_font_themes = {}

        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame is not None:
                for paragraph in shape.text_frame.paragraphs:
                    paragraph_text = ""
                    for run in paragraph.runs:
                        paragraph_text += run.text

                        # Check font name or fallback to theme-based font
                        font_name = run.font.name or "Arial Black or Arial (Body)"

                        # Handle cases where font might be from a theme (e.g., Body or Heading fonts)
                        if font_name in ["Body", "Heading"]:
                            font_name = get_font_from_theme(prs, font_name)

                        # Add font theme information only if there is non-empty text
                        if paragraph_text.strip():
                            slide_font_themes[paragraph_text.strip()] = font_name

        slides_font_themes[slide_num] = slide_font_themes

    return slides_font_themes

def select_file():
    # Open file dialog and allow the user to select a .pptx file
    root = tk.Tk()
    root.withdraw()  # Hide the main tkinter window
    file_path = filedialog.askopenfilename(
        filetypes=[("PowerPoint files", "*.pptx")],
        title="Select a PowerPoint File"
    )
    return file_path

if __name__ == "__main__":
    file_path = select_file()  # Select the file using a dialog
    if file_path:  # If a file was selected
        print(f"Analyzing file: {file_path.split('/')[-1]}")  # Display the file name
        slides_font_themes = check_font_theme_pptx(file_path)

        for slide_num, font_themes in slides_font_themes.items():
            print(f"\nSlide {slide_num}:")
            for text, font_name in font_themes.items():
                print(f"  Text: '{text}' -> Font Theme: {font_name}")
    else:
        print("No file selected.")
