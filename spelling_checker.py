import re
import tkinter as tk
from tkinter import filedialog
from pptx import Presentation
from spellchecker import SpellChecker

# Function to clean words by removing surrounding punctuation while keeping contractions intact
def clean_word(word):
    # Keep words with apostrophes intact but remove other surrounding punctuation
    return re.sub(r"^[\"'.,“”‘’\(\)\[\]]+|[\"'.,“”‘’\(\)\[\]:;?]+$", "", word)

# Function to check spelling mistakes with edge case handling
def check_spelling_pptx(file_path):
    prs = Presentation(file_path)
    spell = SpellChecker()  # Initialize the spell checker
    slides_spelling_mistakes = {}

    # Define patterns to ignore
    contractions = {"we're", "i'm", "they're", "we've", "it's", "i've", "don't", "can't", "today's"}
    ignore_words = {"we're", "i'm", "they're", "we've", "it's", "i've", "don't", "can't", "today's", "CAEs", "caes", "CAE", "cae", "iia", "Gartner", "gartner", "IIA", "DA", "Don't", "don't", "Can't", "Source:", "source:", "etc", "lenovo"}
    # Enhanced regex pattern to ignore words with certain characters and numbers
    ignore_patterns = re.compile(r"^[A-Z0-9()%'“”.,:;?\[\]\(\)]+$")  # Ignore all-caps, numbers, punctuation, brackets

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_mistakes = {}

        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame is not None:
                for paragraph in shape.text_frame.paragraphs:
                    paragraph_text = ""
                    for run in paragraph.runs:
                        paragraph_text += run.text  # Collect the full text in the paragraph

                    words = paragraph_text.split()

                    # Clean words before checking them
                    cleaned_words = [clean_word(word) for word in words]

                    # Ignore words with apostrophes automatically
                    misspelled = spell.unknown([word for word in cleaned_words if "'" not in word])

                    # Filter out ignored patterns, contractions, and predefined ignored words
                    misspelled_filtered = [
                        word for word in misspelled
                        if word.lower() not in contractions and word not in ignore_words and not ignore_patterns.match(word)
                    ]

                    if misspelled_filtered:  # If any filtered misspelled words were found
                        slide_mistakes[paragraph_text.strip()] = misspelled_filtered

        if slide_mistakes:
            slides_spelling_mistakes[slide_num] = slide_mistakes

    return slides_spelling_mistakes

# Function to select a file
def select_file():
    root = tk.Tk()
    root.withdraw()  # Hide the main tkinter window
    file_path = filedialog.askopenfilename(
        filetypes=[("PowerPoint files", "*.pptx")],
        title="Select a PowerPoint File"
    )
    return file_path

# Main Execution
if __name__ == "__main__":
    file_path = select_file()  # Select the file using a dialog
    if file_path:  # If a file was selected
        print(f"Analyzing file: {file_path.split('/')[-1]}")  # Display the file name
        slides_spelling_mistakes = check_spelling_pptx(file_path)

        if slides_spelling_mistakes:
            for slide_num, mistakes in slides_spelling_mistakes.items():
                print(f"\nSlide {slide_num} has spelling mistakes:")
                for text, misspelled_words in mistakes.items():
                    print(f"  Text: '{text}' -> Misspelled words: {misspelled_words}")
        else:
            print("\nNo spelling mistakes found!")
    else:
        print("No file selected.")
